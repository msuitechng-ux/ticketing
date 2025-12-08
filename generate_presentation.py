#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for Graduation Guest Ticketing System
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    # Style title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add content
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(12)

    return slide

def add_two_column_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    """Add a two-column content slide"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    # Style title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Remove default text box
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_shape:
            sp = shape.element
            sp.getparent().remove(sp)

    # Add left column
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.25), Inches(4.5))
    left_frame = left.text_frame
    left_frame.word_wrap = True

    if left_title:
        p = left_frame.paragraphs[0]
        p.text = left_title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(12)
    else:
        left_frame.clear()

    for item in left_items:
        p = left_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    # Add right column
    right = slide.shapes.add_textbox(Inches(5.25), Inches(1.8), Inches(4.25), Inches(4.5))
    right_frame = right.text_frame
    right_frame.word_wrap = True

    if right_title:
        p = right_frame.paragraphs[0]
        p.text = right_title
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 51, 102)
        p.space_after = Pt(12)
    else:
        right_frame.clear()

    for item in right_items:
        p = right_frame.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "Graduation Guest Ticketing System",
        "Northumbria University London Campus\nLD7206 Module Project"
    )

    # Slide 2: System Overview
    add_content_slide(prs, "System Overview", [
        "Purpose: Manage graduation ceremony guest tickets from allocation to entry verification",
        "Designed for: Northumbria University London Campus graduation ceremonies",
        "Stakeholders: Students (Graduates), Administrators, Security Staff, Registry",
        "Key Functions: Ticket allocation, request management, QR verification, fraud detection",
        "Complete lifecycle management from registration through entry logging"
    ])

    # Slide 3: Tech Stack - Backend & Frontend
    add_two_column_slide(
        prs,
        "Modern Tech Stack",
        [
            "Laravel 12 (PHP 8.3)",
            "Laravel Fortify - Authentication & 2FA",
            "Spatie Permissions - RBAC",
            "DomPDF & FPDF - PDF generation",
            "Endroid QR Code - QR generation",
            "Laravel Wayfinder - Type-safe routing",
            "Eloquent ORM - Database management"
        ],
        [
            "Vue 3 with TypeScript",
            "Inertia.js v2 - SSR framework",
            "TailwindCSS v4 - Styling",
            "PrimeVue 4.5 - UI components",
            "Chart.js 4.5 - Data visualization",
            "Vite 7 - Fast module bundling",
            "ESLint & Prettier - Code quality"
        ],
        "Backend",
        "Frontend"
    )

    # Slide 4: Core Data Models
    add_content_slide(prs, "Core Data Models & Relationships", [
        "User: Multi-role accounts (Admin, Student, Security Staff, Registry)",
        "Graduate: Student information, degree level, faculty, department",
        "Ceremony: Event details, capacity, deadlines, ticket allocation rules",
        "Ticket: QR codes, status tracking, type (Base/Extra/Redistributed)",
        "TicketRequest: Student requests for additional tickets with approval workflow",
        "EntryLog: Complete audit trail of all entry attempts with fraud detection"
    ])

    # Slide 5: Key Features - Part 1
    add_content_slide(prs, "Key Features - Administration", [
        "Ceremony Management: Create ceremonies with capacity, deadlines, ticket limits",
        "Graduate Management: Manual entry or bulk CSV import with auto ticket allocation",
        "Ticket Allocation: Automatic base ticket assignment on graduate registration",
        "Request Processing: Approve, deny, waitlist, or partially approve ticket requests",
        "Batch Operations: Process multiple requests simultaneously for efficiency",
        "Analytics Dashboard: Real-time metrics, charts, utilization rates"
    ])

    # Slide 6: Key Features - Part 2
    add_content_slide(prs, "Key Features - Student & Security", [
        "Student Portal: View tickets, submit requests, download QR codes as PDF",
        "QR Code Generation: Unique codes with cryptographic checksums",
        "Entry Verification: Real-time QR scanning with fraud detection",
        "Ticket Downloads: Individual PDFs or bulk ZIP file downloads",
        "Request Tracking: Real-time status updates with admin feedback",
        "Entry Logging: Complete audit trail of all verification attempts"
    ])

    # Slide 7: User Roles & Permissions
    add_two_column_slide(
        prs,
        "Role-Based Access Control",
        [
            "Admin: Full system access, manage ceremonies, graduates, tickets",
            "Registry: Bulk import graduates, manage student data",
            "Event Staff: View analytics and ceremony reports"
        ],
        [
            "Student: View own tickets/requests, submit requests, download QRs",
            "Security Staff: Scan/verify tickets, view entry logs, review fraud cases"
        ],
        "Administrative Roles",
        "Operational Roles"
    )

    # Slide 8: Ticket Lifecycle
    add_content_slide(prs, "Complete Ticket Lifecycle", [
        "1. Graduate Registration → Base tickets automatically allocated (Active)",
        "2. Student Submits Extra Request → Admin receives notification",
        "3. Admin Reviews → Approve/Deny/Waitlist/Partially Approve",
        "4. Approved Tickets Created → QR codes generated with checksums",
        "5. Student Downloads → PDF with embedded QR code",
        "6. Entry Verification → Security staff scans QR at ceremony",
        "7. Status Updated → Ticket marked as Used, entry logged",
        "8. Redistribution → Unused tickets reallocated to waitlist if deadline passed"
    ])

    # Slide 9: Advanced QR Code System
    add_content_slide(prs, "Sophisticated QR Code Technology", [
        "Cryptographic Validation: Each QR contains SHA256 checksum",
        "Data Structure: JSON with ticket_id, code, ceremony_id, graduate_id, type, checksum",
        "Base64 Encoding: Optimized for PDF embedding and scanning",
        "PDF Generation: Professional ticket design with embedded QR",
        "Bulk Operations: Generate and download all tickets as ZIP",
        "Regeneration Capability: Lost/corrupted codes can be regenerated",
        "Multi-layer Validation: Checksum, status, duplicate, and fraud detection"
    ])

    # Slide 10: Ticket Request Workflow
    add_content_slide(prs, "Intelligent Request Management", [
        "Student Submission: Request additional tickets with justification/reason",
        "Admin Review Queue: Prioritized by submission date",
        "Flexible Decisions: Approve exact amount, partial amount, deny, or waitlist",
        "Admin Notes: Provide feedback and reasoning to students",
        "Partial Approval: Can approve fewer tickets than requested",
        "Waitlist System: Automatic redistribution when tickets become available",
        "Deadline Enforcement: Respects per-ceremony request and redistribution deadlines",
        "Duplicate Prevention: Only one active request per student at a time"
    ])

    # Slide 11: Security & Fraud Detection
    add_content_slide(prs, "Multi-Layer Security Features", [
        "Checksum Validation: Detects tampered or forged QR codes",
        "Duplicate Detection: Prevents scanning the same ticket twice",
        "Status Verification: Checks ticket validity, expiration, and active status",
        "Entry Logging: Complete audit trail with timestamp, scanner, device, location",
        "Fraud Case Management: Dedicated interface to review suspicious entries",
        "Real-time Alerts: Instant notification of fraud attempts",
        "Device Tracking: Logs scanner device information for accountability",
        "Role-Based Authorization: Middleware-enforced access control on all routes"
    ])

    # Slide 12: Analytics & Reporting
    add_content_slide(prs, "Comprehensive Analytics Dashboard", [
        "Real-time Metrics: Total ceremonies, graduates, tickets, utilization rates",
        "Interactive Charts: Faculty distribution, degree levels, ticket allocation vs usage",
        "Request Analytics: Status breakdown (pending, approved, denied, waitlisted)",
        "Entry Statistics: Successful scans, fraud attempts, duplicate attempts",
        "Ceremony Analytics: Per-ceremony performance and capacity tracking",
        "Upcoming Events: Next 5 ceremonies with graduate counts",
        "Recent Activity: Latest requests and entry logs for quick monitoring",
        "Staff Performance: Track which security staff scanned which tickets"
    ])

    # Slide 13: System Architecture
    add_content_slide(prs, "Clean Architecture Patterns", [
        "Service Layer: TicketAllocationService, TicketRequestService, EntryVerificationService, QrCodeService",
        "Resource Controllers: RESTful endpoints for ceremonies, graduates, tickets",
        "Form Request Validation: Server-side validation with custom error messages",
        "Eloquent Relationships: Optimized queries with eager loading to prevent N+1",
        "Database Transactions: Ensures data integrity on critical operations",
        "Event-Driven: Laravel events for ticket creation, status changes",
        "API Resources: Structured JSON responses for frontend consumption",
        "Repository Pattern: Clean separation of data access logic"
    ])

    # Slide 14: Standout Features
    add_content_slide(prs, "What Makes This System Unique", [
        "✓ Partial Request Approval: Approve fewer tickets than requested intelligently",
        "✓ Automatic Redistribution: Waitlisted students automatically receive unused tickets",
        "✓ Cryptographic QR Security: Industry-standard SHA256 checksums",
        "✓ Comprehensive Audit Trail: Every action logged with complete traceability",
        "✓ Bulk CSV Import: Import hundreds of graduates in seconds",
        "✓ Fraud Detection AI: Multi-layer validation catches tampering attempts",
        "✓ Real-time Analytics: Live dashboard updates without page refresh",
        "✓ Mobile-Responsive: Works seamlessly on tablets and phones for scanning"
    ])

    # Slide 15: Database Schema Highlights
    add_two_column_slide(
        prs,
        "Database Schema Design",
        [
            "Tickets: ticket_code, qr_code_path, status, is_scanned, scanned_at, ticket_type",
            "Graduates: student_id, degree_level, faculty, department",
            "Ceremonies: date, venue, capacity, base_tickets_per_graduate"
        ],
        [
            "TicketRequests: status, requested_qty, approved_qty, reviewed_by, admin_notes",
            "EntryLogs: verification_status, scanned_by, entry_point, device_info",
            "Users: Multi-role with Spatie permissions"
        ],
        "Primary Tables",
        "Supporting Tables"
    )

    # Slide 16: Batch Operations
    add_content_slide(prs, "Efficient Batch Processing", [
        "Bulk Graduate Import: CSV upload with validation and error reporting",
        "Batch Request Processing: Review and approve multiple requests at once",
        "Batch Ticket Generation: Create tickets for all approved requests simultaneously",
        "Bulk QR Downloads: ZIP file export of all tickets for a ceremony",
        "Batch Fraud Review: Review multiple suspicious entries together",
        "Automatic Redistribution: Process entire waitlist when tickets available",
        "Transaction Safety: All batch operations wrapped in database transactions"
    ])

    # Slide 17: Testing & Quality Assurance
    add_content_slide(prs, "Quality Assurance", [
        "Pest Testing Framework: Modern PHP testing with expressive syntax",
        "Feature Tests: Complete user flow testing from registration to entry",
        "Unit Tests: Individual service and model method testing",
        "Browser Tests: End-to-end testing with real browsers (Chrome, Firefox)",
        "Laravel Pint: Automatic code formatting and style enforcement",
        "ESLint & Prettier: Frontend code quality and consistency",
        "Form Request Validation: Server-side validation on all inputs",
        "Security Testing: QR tampering, duplicate attempts, fraud scenarios"
    ])

    # Slide 18: Deployment & Infrastructure
    add_content_slide(prs, "Production Ready", [
        "Laravel Sail: Docker-based development environment for consistency",
        "Vite HMR: Hot module replacement for fast frontend development",
        "Database: SQLite (dev), MySQL/PostgreSQL ready for production",
        "File Storage: Local (dev), S3-compatible storage for production",
        "Queue System: Database queues ready for Redis/SQS in production",
        "Caching: Application, route, and config caching for performance",
        "Environment Config: Proper separation of dev/staging/production settings",
        "Git Workflow: Feature branches with pull request reviews"
    ])

    # Slide 19: User Experience Highlights
    add_content_slide(prs, "Superior User Experience", [
        "Responsive Design: Works perfectly on desktop, tablet, and mobile devices",
        "Real-time Updates: Instant feedback on all actions without page reloads",
        "Intuitive Navigation: Clean, logical flow through all user journeys",
        "Professional PDFs: Well-designed ticket downloads with QR codes",
        "Error Handling: Clear, helpful error messages guide users",
        "Loading States: Skeleton loaders and spinners for better perceived performance",
        "Dark Mode Support: TailwindCSS dark mode for reduced eye strain",
        "Accessibility: Semantic HTML and ARIA labels for screen readers"
    ])

    # Slide 20: Project Statistics
    add_content_slide(prs, "Project Metrics", [
        "Backend: 6 Eloquent Models, 8 Controllers, 4 Service Classes",
        "Frontend: 20+ Vue Components, 15+ Page Components",
        "Database: 10 Migration Files, complete relational schema",
        "Validation: 7 Form Request Classes with custom error messages",
        "Testing: Feature tests, unit tests, browser tests with Pest v4",
        "Dependencies: 25+ Composer packages, 20+ NPM packages",
        "Code Quality: Laravel Pint formatting, ESLint + Prettier for JS",
        "Documentation: Comprehensive code comments and PHPDoc blocks"
    ])

    # Slide 21: Future Enhancements
    add_content_slide(prs, "Potential Future Features", [
        "SMS/Email Notifications: Auto-notify students of ticket approvals and ceremony reminders",
        "Mobile App: Native iOS/Android apps for ticket management and scanning",
        "Digital Wallet Integration: Add tickets to Apple Wallet / Google Pay",
        "Guest Information Collection: Capture guest names and contact details",
        "Seating Assignments: Allocate specific seats based on ticket codes",
        "Real-time Capacity Dashboard: Live view of ceremony attendance",
        "Multi-Language Support: Interface translation for international students",
        "API for Third-party Integration: Allow integration with university systems"
    ])

    # Slide 22: Technical Challenges Solved
    add_content_slide(prs, "Complex Problems Addressed", [
        "Challenge: Prevent QR code forgery → Solution: SHA256 cryptographic checksums",
        "Challenge: Fair ticket distribution → Solution: Waitlist with automatic redistribution",
        "Challenge: Duplicate entry attempts → Solution: Entry logging with duplicate detection",
        "Challenge: Scalability for large ceremonies → Solution: Efficient queries with eager loading",
        "Challenge: User-friendly scanning → Solution: Real-time validation with clear feedback",
        "Challenge: Data integrity → Solution: Database transactions and foreign key constraints",
        "Challenge: Request deadline management → Solution: Automated deadline enforcement",
        "Challenge: Audit compliance → Solution: Complete entry logs with immutable records"
    ])

    # Slide 23: Business Value
    add_content_slide(prs, "Business Impact & Benefits", [
        "Efficiency: Automated allocation saves hours of manual work",
        "Accuracy: Eliminates errors from manual ticket distribution",
        "Security: Fraud detection protects ceremony integrity",
        "Transparency: Complete audit trail for compliance and disputes",
        "Scalability: Handles hundreds of graduates and thousands of tickets",
        "Cost Savings: Reduces staff time and eliminates paper tickets",
        "Student Satisfaction: Clear process, instant feedback, easy downloads",
        "Data Insights: Analytics inform future ceremony planning"
    ])

    # Slide 24: System Demonstration Flow
    add_content_slide(prs, "Live Demo Walkthrough", [
        "1. Admin creates ceremony and sets ticket allocation rules",
        "2. Registry imports graduates via CSV bulk upload",
        "3. System automatically allocates base tickets to all graduates",
        "4. Student logs in and requests additional tickets with reason",
        "5. Admin reviews request and partially approves (2 of 3 requested)",
        "6. Student downloads PDF tickets with QR codes",
        "7. Security staff scans QR at ceremony entrance",
        "8. System validates, logs entry, and displays success/fraud alert",
        "9. Admin views analytics dashboard with real-time statistics"
    ])

    # Slide 25: Learning Outcomes
    add_content_slide(prs, "Key Learning & Skills Developed", [
        "Full-stack development with modern Laravel and Vue ecosystem",
        "Database design with complex relationships and constraints",
        "RESTful API design and implementation",
        "Authentication, authorization, and role-based access control",
        "QR code generation and cryptographic validation",
        "PDF generation and document management",
        "Real-time data visualization with Chart.js",
        "Testing strategies: unit, feature, and browser tests",
        "DevOps: Docker, version control, CI/CD readiness"
    ])

    # Slide 26: Technologies Mastered
    add_two_column_slide(
        prs,
        "Technical Skills Demonstrated",
        [
            "Laravel 12 framework",
            "Eloquent ORM & Query Builder",
            "Inertia.js SSR architecture",
            "Vue 3 Composition API",
            "TypeScript integration",
            "Spatie Laravel Permission",
            "DomPDF & FPDF libraries"
        ],
        [
            "TailwindCSS v4 utility classes",
            "Chart.js data visualization",
            "QR code generation & validation",
            "Pest v4 testing framework",
            "Laravel Wayfinder routing",
            "Vite build optimization",
            "Docker containerization"
        ],
        "Backend Mastery",
        "Frontend Excellence"
    )

    # Slide 27: Code Quality Standards
    add_content_slide(prs, "Best Practices Implemented", [
        "PSR-12 Code Style: Enforced by Laravel Pint automatically",
        "SOLID Principles: Service layer separates business logic from controllers",
        "DRY (Don't Repeat Yourself): Reusable components and services",
        "Security First: OWASP top 10 vulnerabilities addressed",
        "Type Safety: TypeScript and PHP type hints throughout",
        "Clean Code: Descriptive naming, single responsibility functions",
        "Documentation: PHPDoc blocks and inline comments where needed",
        "Version Control: Semantic commits with clear messages"
    ])

    # Slide 28: Conclusion
    add_content_slide(prs, "Project Summary", [
        "Successfully delivered a production-grade ticketing system",
        "Comprehensive feature set covering entire ticket lifecycle",
        "Modern tech stack with Laravel 12 and Vue 3",
        "Advanced security with QR validation and fraud detection",
        "Scalable architecture ready for thousands of users",
        "Clean, maintainable code following best practices",
        "Complete testing coverage for reliability",
        "Valuable learning experience in full-stack development"
    ])

    # Slide 29: Thank You
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Thank You!"

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add subtitle
    left = Inches(1)
    top = Inches(3.8)
    width = Inches(8)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Graduation Guest Ticketing System\nNorthumbria University London Campus"

    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = RGBColor(100, 100, 100)

    # Save presentation
    output_file = "/home/user/ticketing/Graduation_Ticketing_System_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_presentation()
